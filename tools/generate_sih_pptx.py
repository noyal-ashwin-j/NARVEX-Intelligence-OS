import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

def add_transition(slide, transition_type="fade"):
    """Adds native PowerPoint transition animation (e.g. fade, push, wipe)"""
    try:
        transition_xml = f'<p:transition {nsdecls("p")}><p:{transition_type}/></p:transition>'
        slide._element.append(parse_xml(transition_xml))
    except Exception as e:
        print(f"Warning: Could not set transition: {e}")

def create_styled_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Tactical Sovereign Intelligence Color Palette
    DARK_NAVY = RGBColor(9, 14, 26)       # #090E1A
    MID_NAVY = RGBColor(15, 23, 42)       # #0F172A
    WHITE = RGBColor(255, 255, 255)       # Clean white
    CARD_BG = RGBColor(255, 255, 255)     # Clean white card
    CARD_DARK = RGBColor(17, 24, 39)      # #111827
    CYAN_NEON = RGBColor(34, 211, 238)    # #22D3EE
    CYAN_DEEP = RGBColor(14, 116, 144)    # #0E7490
    AMBER_ALERT = RGBColor(245, 158, 11)  # #F59E0B
    ROSE_CRITICAL = RGBColor(239, 68, 68) # #EF4444
    EMERALD_PASS = RGBColor(16, 185, 129) # #10B981
    PURPLE_AI = RGBColor(147, 51, 234)    # #9333EA
    TEXT_DARK = RGBColor(30, 41, 59)
    TEXT_LIGHT = RGBColor(248, 250, 252)
    TEXT_MUTED = RGBColor(100, 116, 139)
    BORDER_LIGHT = RGBColor(226, 232, 240)
    BORDER_CYAN = RGBColor(34, 211, 238)

    blank_layout = prs.slide_layouts[6]

    def add_top_header(slide, title_text, subtitle_text="SMART INDIA HACKATHON 2026"):
        # Header Base Bar
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = DARK_NAVY
        top_bar.line.fill.background()

        # Glowing Cyan Accent Stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.06))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = CYAN_NEON
        stripe.line.fill.background()

        # Subtitle / Category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(8), Inches(0.3))
        cat_tf = cat_box.text_frame
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = subtitle_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN_NEON

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(10), Inches(0.65))
        title_tf = title_box.text_frame
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(21)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_LIGHT

        # Footer Bar
        foot_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.12), Inches(13.333), Inches(0.38))
        foot_bar.fill.solid()
        foot_bar.fill.fore_color.rgb = DARK_NAVY
        foot_bar.line.fill.background()

        foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.16), Inches(11.733), Inches(0.3))
        foot_tf = foot_box.text_frame
        p_foot = foot_tf.paragraphs[0]
        p_foot.text = "NARVEX — Sovereign Spatial-Temporal Narcotics Intelligence & Predictive OS | @SIH Idea Submission"
        p_foot.font.size = Pt(9)
        p_foot.font.color.rgb = RGBColor(148, 163, 184)

    # -------------------------------------------------------------
    # SLIDE 1: TITLE PAGE (With Fade Transition & Hero Cards)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_transition(slide1, "fade")

    # Background
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(241, 245, 249)
    bg1.line.fill.background()

    add_top_header(slide1, "TITLE PAGE", "SMART INDIA HACKATHON 2026")

    # Left Container (Metadata Card)
    left_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.3))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = BORDER_LIGHT

    # Metadata Top Header Strip
    lt_strip = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(7.2), Inches(0.55))
    lt_strip.fill.solid()
    lt_strip.fill.fore_color.rgb = DARK_NAVY
    lt_strip.line.fill.background()

    lt_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.58), Inches(6.8), Inches(0.4))
    lt_tf = lt_box.text_frame
    p_lt = lt_tf.paragraphs[0]
    p_lt.text = "PROBLEM STATEMENT & TEAM REGISTRATION"
    p_lt.font.size = Pt(12)
    p_lt.font.bold = True
    p_lt.font.color.rgb = CYAN_NEON

    l_content = slide1.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(6.6), Inches(4.4))
    l_tf = l_content.text_frame
    l_tf.word_wrap = True

    items = [
        ("Problem Statement ID", "SIH2026-NCB-01 (or Allocated ID)"),
        ("Problem Statement Title", "Unified Spatial-Temporal Intelligence, Cross-Agency Signal Correlation & Predictive Narcotics Corridor Tracking Platform"),
        ("Theme", "Security & Surveillance / Smart Governance / AI Decision Support"),
        ("PS Category", "Software (Enterprise / Air-Gapped Sovereign Intelligence OS)"),
        ("Team ID", "[Insert Team ID]"),
        ("Team Name", "[Insert Registered Team Name]")
    ]

    for i, (label, val) in enumerate(items):
        p = l_tf.add_paragraph() if i > 0 else l_tf.paragraphs[0]
        p.space_after = Pt(10)
        r1 = p.add_run()
        r1.text = f"• {label}: "
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = DARK_NAVY

        r2 = p.add_run()
        r2.text = val
        r2.font.bold = False
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_DARK

    # Right Hero Card (Futuristic Cyber Aesthetic)
    right_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.5), Inches(4.333), Inches(5.3))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = DARK_NAVY
    right_card.line.color.rgb = CYAN_NEON

    r_box = slide1.shapes.add_textbox(Inches(8.4), Inches(1.8), Inches(3.933), Inches(4.7))
    r_tf = r_box.text_frame
    r_tf.word_wrap = True

    p1 = r_tf.paragraphs[0]
    p1.text = "NARVEX"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_NEON
    p1.alignment = PP_ALIGN.CENTER

    p2 = r_tf.add_paragraph()
    p2.text = "Sovereign Intelligence OS"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(14)

    p3 = r_tf.add_paragraph()
    p3.text = "JARVIS : IRON MAN\nNARVEX : LAW ENFORCEMENT"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = AMBER_ALERT
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(14)

    p4 = r_tf.add_paragraph()
    p4.text = "• Cross-Agency Signal Correlation\n• 4-Tier Grounded Corridor Telemetry\n• Two-Horizon Predictive Modeling\n• Tamper-Evident SHA-256 Provenance\n• Zero-Trust Role-Based Access Control"
    p4.font.size = Pt(11)
    p4.font.color.rgb = RGBColor(226, 232, 240)

    # -------------------------------------------------------------
    # SLIDE 2: PROPOSED SOLUTION (With Wipe Transition)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_transition(slide2, "wipe")
    add_top_header(slide2, "PROPOSED SOLUTION (NARVEX INTELLIGENCE PLATFORM)")

    def add_feature_card(slide, x, y, w, h, title, badge_text, badge_color, points, header_color=DARK_NAVY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_LIGHT

        # Top Header Strip
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.55))
        tag.fill.solid()
        tag.fill.fore_color.rgb = header_color
        tag.line.fill.background()

        tbox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(w - 0.3), Inches(0.4))
        ttf = tbox.text_frame
        p = ttf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Badge Pill
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(y + 0.65), Inches(w - 0.3), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_color
        badge.line.fill.background()

        bbox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.65), Inches(w - 0.3), Inches(0.35))
        btf = bbox.text_frame
        bp = btf.paragraphs[0]
        bp.text = badge_text
        bp.font.size = Pt(9.5)
        bp.font.bold = True
        bp.font.color.rgb = WHITE
        bp.alignment = PP_ALIGN.CENTER

        cbox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.05), Inches(w - 0.3), Inches(h - 1.15))
        ctf = cbox.text_frame
        ctf.word_wrap = True

        for i, pt_text in enumerate(points):
            p_pt = ctf.add_paragraph() if i > 0 else ctf.paragraphs[0]
            p_pt.space_after = Pt(8)
            p_pt.text = f"• {pt_text}"
            p_pt.font.size = Pt(10.5)
            p_pt.font.color.rgb = TEXT_DARK

    s2_p1 = [
        "Continuous Intelligence OS: Not a static dashboard or simple complaint form; continuously recalculates risk.",
        "Universal Drag-and-Drop Ingestion: Accepts PDF FIRs, Seizure logs, Toll ANPR scans, Citizen tips & News.",
        "Two-Horizon Analytics: Differentiates Immediate Attention (Velocity ≥ 1.8x) from 30-Day ML Risk Horizon."
    ]
    add_feature_card(slide2, 0.8, 1.5, 3.65, 5.3, "1. PROPOSED SOLUTION", "SELF-UPDATING INTELLIGENCE OS", DARK_NAVY, s2_p1, DARK_NAVY)

    s2_p2 = [
        "Root Cause: Disconnected crimes near colleges/borders are handled in silos until violent crime occurs.",
        "Cross-Source Correlation: Fuses data from State Police, Toll ANPR, Customs, and Citizen tips in 30D windows.",
        "Early Cluster Detection: Discovers micro-clusters and cross-agency convergence weeks before crisis escalation."
    ]
    add_feature_card(slide2, 4.84, 1.5, 3.65, 5.3, "2. SOLVING SIGNAL FAILURE", "CROSS-AGENCY CORRELATION", CYAN_DEEP, s2_p2, CYAN_DEEP)

    s2_p3 = [
        "4-Tier Corridor Evidence: Classifies routes as Observed, Associated, Potential, or Unknown.",
        "First-Time Safeguard: Zero-history areas flagged as Needs Verification rather than false High Risk.",
        "Informant Shield: Automated PII stripping + anonymous cryptographic tracking tokens (NARC-xxxx).",
        "Explainable SHA-256 Provenance: 112+ block audit ledger answering 'Why is this flagged?'."
    ]
    add_feature_card(slide2, 8.88, 1.5, 3.65, 5.3, "3. INNOVATION & UNIQUENESS", "CRYPTOGRAPHIC & EXPLAINABLE", PURPLE_AI, s2_p3, PURPLE_AI)

    # -------------------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH (With Push Transition)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_transition(slide3, "push")
    add_top_header(slide3, "TECHNICAL APPROACH & SYSTEM ARCHITECTURE")

    # Flow Banner
    flow_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.2))
    flow_card.fill.solid()
    flow_card.fill.fore_color.rgb = DARK_NAVY
    flow_card.line.color.rgb = CYAN_NEON

    fbox = slide3.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.333), Inches(2.0))
    ftf = fbox.text_frame
    ftf.word_wrap = True

    p_fl_title = ftf.paragraphs[0]
    p_fl_title.text = "⚡ AUTOMATED END-TO-END INTELLIGENCE PIPELINE"
    p_fl_title.font.size = Pt(12)
    p_fl_title.font.bold = True
    p_fl_title.font.color.rgb = CYAN_NEON
    p_fl_title.space_after = Pt(4)

    pipeline_steps = (
        "1. RAW MULTI-SOURCE FEEDS (PDF FIRs • Seizures • Toll ANPR • Citizen Reports • Public News)\n"
        "   ↓ Automated PII Redaction + Tesseract OCR Parser + Geospatial Normalization\n"
        "2. NORMALIZED MYSQL narvex LEDGER (Historical 30D/90D Baselines & Cryptographic Blocks)\n"
        "   ↓ Spatio-Temporal Correlation Engine & 4-Tier Corridor Evidence Grounding\n"
        "3. CALIBRATED ML FORECASTING (Regularized Logistic Regression, Temperature T=1.6, P ∈ [0.15, 0.88])\n"
        "   ↓ Real-Time Server-Sent Events (SSE Mesh Broadcast)\n"
        "4. TACTICAL COMMAND HUD (MapLibre 3D Globe Arcs • Leaflet 2D GIS • AI Voice Assistant)"
    )
    p_fl_body = ftf.add_paragraph()
    p_fl_body.text = pipeline_steps
    p_fl_body.font.size = Pt(9.5)
    p_fl_body.font.color.rgb = RGBColor(226, 232, 240)

    # 3 Tech Stack Cards
    t1 = [
        "React 18 & Vite: High-performance command canvas.",
        "MapLibre GL 3D: WebGL great-circle interstate arcs.",
        "Leaflet 2D GIS: District street-level surveillance.",
        "TailwindCSS & Lucide: Tactical war-room UI."
    ]
    add_feature_card(slide3, 0.8, 3.9, 3.65, 2.9, "FRONTEND & VISUALIZATION", "REACT 18 • MAPLIBRE 3D", DARK_NAVY, t1, DARK_NAVY)

    t2 = [
        "Node.js & Express: Parameterized API gateway.",
        "SSE Real-Time Mesh: Sub-second live telemetry.",
        "Regularized ML Model: Calibrated 5-feature vector.",
        "NLP Intent Parser: Central Voice Assistant tool bus."
    ]
    add_feature_card(slide3, 4.84, 3.9, 3.65, 2.9, "BACKEND & AI ENGINE", "NODE.JS • SSE • ML REGRESSION", CYAN_DEEP, t2, CYAN_DEEP)

    t3 = [
        "MySQL 8.0: 38-district intelligence schema.",
        "Bcrypt Hashing: Strict salted password verification.",
        "RFC 6238 TOTP: Time-window Multi-Factor Auth.",
        "SHA-256 Ledger: 112+ block audit hash chain."
    ]
    add_feature_card(slide3, 8.88, 3.9, 3.65, 2.9, "SECURITY & LEDGER", "ZERO-TRUST • SHA-256 HASH", PURPLE_AI, t3, PURPLE_AI)

    # -------------------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY (With Wipe Transition)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_transition(slide4, "wipe")
    add_top_header(slide4, "FEASIBILITY, VIABILITY & RISK MITIGATION")

    f_left = [
        "100% Working Prototype Already Built: 38 Tamil Nadu district nodes populated, 2,400+ normalized records, 10/10 automated tests passing.",
        "Zero Officer Training Curve: Directly ingests standard CCTNS FIR exports and ANPR toll logs with drag-and-drop ease.",
        "Government Policy Alignment: Full compliance with National Data Governance Framework and IT Act informant privacy mandates.",
        "Hardware Independence: Runs on standard state command workstations with zero expensive proprietary hardware lock-in."
    ]
    add_feature_card(slide4, 0.8, 1.5, 5.65, 5.3, "FEASIBILITY & OPERATIONAL VIABILITY", "PRODUCTION READY & COMPLIANT", DARK_NAVY, f_left, DARK_NAVY)

    r_right = [
        "Challenge 1 (Unstructured / Scanned FIRs): Solved via multi-pass OCR (Tesseract) + regex entity extraction with human verification queue.",
        "Challenge 2 (Insider Threat / Data Leakage): Solved via Server-Side Zero-Trust District Scoping returning HTTP 403 & SIEM alerts.",
        "Challenge 3 (Air-Gapped Police Networks): Solved via bundled offline PMTiles vector maps and local ONNX inference nodes without cloud APIs.",
        "Challenge 4 (Informant Retaliation Risk): Solved via automated PII stripping and zero-knowledge tracking tokens (NARC-xxxx)."
    ]
    add_feature_card(slide4, 6.88, 1.5, 5.65, 5.3, "CHALLENGES & MITIGATION STRATEGY", "ZERO-TRUST SECURITY ENVELOPE", ROSE_CRITICAL, r_right, ROSE_CRITICAL)

    # -------------------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS (With Fade Transition)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_transition(slide5, "fade")
    add_top_header(slide5, "IMPACT AND SOCIETAL / ADMINISTRATIVE BENEFITS")

    i1 = [
        "85% Faster Early Warning: Detects supply corridor shifts weeks before retrospective monthly bulletins.",
        "Automated Feed Processing: Parses multi-format documents in <3 seconds without data entry fatigue.",
        "1-Click Court-Ready Dossiers: Generates tamper-evident executive briefings with full cryptographic evidence lineage.",
        "Optimized Resource Allocation: Directs highway patrols and checkpost intensity to high-probability corridors."
    ]
    add_feature_card(slide5, 0.8, 1.5, 5.65, 5.3, "LAW ENFORCEMENT & ADMINISTRATIVE IMPACT", "85% FASTER SURVEILLANCE", CYAN_DEEP, i1, CYAN_DEEP)

    i2 = [
        "Preventing Violent Crime Escalation: Correlates micro-seizures around colleges before syndicates establish violent territorial control.",
        "Informant & Whistleblower Safety: Eliminates identity leakages, fostering higher citizen reporting trust.",
        "Inter-Agency Synergy: Eliminates duplicated investigations between State Police, NCB, Toll Checkposts, and Health agencies.",
        "Economic Impact: Drastically reduces illicit cash flows financing organized crime syndicates."
    ]
    add_feature_card(slide5, 6.88, 1.5, 5.65, 5.3, "SOCIETAL & ECONOMIC BENEFITS", "COMMUNITY & YOUTH SAFETY", EMERALD_PASS, i2, EMERALD_PASS)

    # -------------------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES (With Push Transition)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_transition(slide6, "push")
    add_top_header(slide6, "RESEARCH, EMPIRICAL BENCHMARKS & REFERENCES")

    ref_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
    ref_card.fill.solid()
    ref_card.fill.fore_color.rgb = WHITE
    ref_card.line.color.rgb = BORDER_LIGHT

    r_strip = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.55))
    r_strip.fill.solid()
    r_strip.fill.fore_color.rgb = DARK_NAVY
    r_strip.line.fill.background()

    rs_box = slide6.shapes.add_textbox(Inches(1.0), Inches(1.58), Inches(11.333), Inches(0.4))
    rs_tf = rs_box.text_frame
    p_rs = rs_tf.paragraphs[0]
    p_rs.text = "EMPIRICAL GROUNDING & GOVERNMENT STANDARDS"
    p_rs.font.size = Pt(12)
    p_rs.font.bold = True
    p_rs.font.color.rgb = CYAN_NEON

    rbox = slide6.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(11.133), Inches(4.3))
    rtf = rbox.text_frame
    rtf.word_wrap = True

    references = [
        ("1. Narcotics Control Bureau (NCB) Intelligence Reports (2025–2026)", "Seizure pattern analysis, multi-modal trafficking trends (long-distance rail, private inter-state buses), and Tamil Nadu maritime corridor telemetry (narcoticsindia.nic.in)."),
        ("2. Ministry of Home Affairs (MHA) MANAS Portal Telemetry", "1.19 Lakh citizen interactions recorded in 2025, validating the critical necessity for automated cross-source signal correlation rather than passive storage."),
        ("3. Smart India Hackathon (SIH) National Problem Track Precedents", "Dual-use chemical diversion tracking (NC049) and cryptocurrency money trail identification in illicit trafficking."),
        ("4. Academic Spatial-Temporal Modeling Research", "Spatial-Temporal Point Processes and Hawkes Self-Exciting Models in Illicit Supply Chain Disruption (IEEE / ACM Intelligence Systems)."),
        ("5. Cryptographic & Security Standards", "NIST FIPS 180-4 (Secure Hash Standard SHA-256), RFC 6238 (TOTP Algorithm), and OWASP Zero-Trust Application Security Guidelines.")
    ]

    for i, (heading, desc) in enumerate(references):
        p_h = rtf.add_paragraph() if i > 0 else rtf.paragraphs[0]
        p_h.space_after = Pt(4)
        r_h = p_h.add_run()
        r_h.text = f"• {heading}: "
        r_h.font.bold = True
        r_h.font.size = Pt(12)
        r_h.font.color.rgb = DARK_NAVY

        r_d = p_h.add_run()
        r_d.text = desc
        r_d.font.bold = False
        r_d.font.size = Pt(11)
        r_d.font.color.rgb = TEXT_DARK

    output_path = "NARVEX_SIH_2026_ANIMATED_PITCH.pptx"
    prs.save(output_path)
    print(f"Enhanced Styled Presentation with Animations saved successfully to {output_path}")

if __name__ == "__main__":
    create_styled_presentation()
